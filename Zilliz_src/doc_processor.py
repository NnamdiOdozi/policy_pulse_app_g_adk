import os
import sys
import datetime
import time
from typing import List, Dict, Any, Optional, Tuple
from tqdm import tqdm
import argparse

# For document processing
import docx
import fitz
import re
from langchain.text_splitter import RecursiveCharacterTextSplitter, CharacterTextSplitter

# Load environment variables
from dotenv import load_dotenv
load_dotenv()
import hashlib  # For SHA-256 hash calculation

from pathlib import Path

# Calculate project root based on file location
current_file = Path(__file__).resolve()
project_root = current_file.parent.parent  # Adjust levels as needed
sys.path.insert(0, str(project_root))

import Utils.client_factory as client_factory

# Configuration
LLM_SUMMARISER_MODEL = "gpt-3.5-turbo"  # Model for generating summaries. can switch up to gpt-4o-mini when doing evals
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=CHUNK_SIZE,
    chunk_overlap=CHUNK_OVERLAP,
    length_function=len,
)   

# Now include the DocumentProcessor class
class DocumentProcessor:
    """Process documents into chunks with rich metadata."""
    
    def __init__(self, openai_client):
        """
         Initialize DocumentProcessor with injected OpenAI client.
        
        Args:
            openai_client: Pre-configured ChatOpenAI client
        """
        self.text_splitter = text_splitter
        
        # Initialize LLM for generating summaries
        self.llm = openai_client
        
        print(f"DocumentProcessor object initialized with LLM model {LLM_SUMMARISER_MODEL}")

    def extract_text_from_file(self, file_path: Path) -> str:
        """
        Extract text content from different file types.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Extracted text content
        """
        file_extension = file_path.suffix.lower()
        
        if file_extension == '.pdf':
            return self._extract_from_pdf(file_path)
        elif file_extension == '.docx':
            return self._extract_from_docx(file_path)
        elif file_extension == '.pptx':
            return self._extract_from_pptx(file_path)
        elif file_extension == '.odp':
            return self._extract_from_odp(file_path)
        elif file_extension in ['.txt', '.md']:
            return self._extract_from_text(file_path)
        else:
            print(f"Unsupported file type: {file_extension}")
            return ""


    def calculate_file_hash(self, file_path: Path) -> str:
        """
        Calculate SHA-256 hash of file content.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Hex string of SHA-256 hash (64 characters)
        """
        sha256_hash = hashlib.sha256()
        
        try:
            with open(file_path, "rb") as f:
                # Read file in chunks to handle large files efficiently
                for byte_block in iter(lambda: f.read(65536), b""):
                    sha256_hash.update(byte_block)
            
            return sha256_hash.hexdigest()
        except Exception as e:
            print(f"Error calculating hash for {file_path}: {e}")
            return ""
        
    def _extract_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint (PPTX) files."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            print(f"Error extracting text from PPTX {file_path}: {e}")
            return ""

    def _extract_from_odp(self, file_path: Path) -> str:
        """Extract text from OpenDocument Presentation (ODP) files."""
        try:
            from odf import text, teletype
            from odf.opendocument import load
            doc = load(file_path)
            text_elements = doc.getElementsByType(text.P)
            return "\n".join([teletype.extractText(element) for element in text_elements])
        except Exception as e:
            print(f"Error extracting text from ODP {file_path}: {e}")
            return ""

    def _extract_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF files with OCR fallback for scanned PDFs."""
        try:
              # PyMuPDF
            
            # Open the PDF
            doc = fitz.open(file_path)
            
            # Check if the PDF has text
            text = ""
            text_found = False
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                page_text = page.get_text()
                
                # If page has more than 10 characters, consider it a text PDF
                if len(page_text.strip()) > 10:
                    text_found = True
                    text += page_text + "\n"
            
            # If no significant text found, it's likely a scanned PDF
            if not text_found:
                print(f"PDF appears to be scanned, applying OCR: {file_path}")
                return self._extract_from_scanned_pdf(file_path)
            
            return text
            
        except Exception as e:
            print(f"Error extracting text from PDF {file_path}: {e}")
            return ""

    def _extract_from_scanned_pdf(self, file_path: Path) -> str:
        """Extract text from scanned PDFs using OCR."""
        try:
            from pdf2image import convert_from_path
            import pytesseract
            
            # Convert PDF to images
            images = convert_from_path(file_path)
            
            # Apply OCR to each image
            text = []
            for i, image in enumerate(images):
                text.append(pytesseract.image_to_string(image))
            
            return "\n".join(text)
            
        except Exception as e:
            print(f"Error extracting text from scanned PDF {file_path}: {e}")
            return ""
    
    def _extract_from_docx(self, file_path: Path) -> str:
        """Extract text from DOCX files."""
        try:
            doc = docx.Document(file_path)
            return "\n\n".join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            print(f"Error extracting text from DOCX {file_path}: {e}")
            return ""
    
    def _extract_from_text(self, file_path: Path) -> str:
        """Extract text from TXT or MD files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            print(f"Error extracting text from text file {file_path}: {e}")
            return ""
    
    def extract_section_title(self, text: str) -> str:
        """
        Extract a section title from the chunk text.
        
        Args:
            text: The chunk text
            
        Returns:
            Extracted section title or default title
        """
        # Look for common section title patterns
        section_patterns = [
            r'^(?:Section|SECTION)\s+\d+[\.:]\s*(.*?)(?:\n|$)',  # Section X: Title
            r'^(?:\d+\.)+\s*(.*?)(?:\n|$)',  # 1.2.3 Title
            r'^[A-Z][A-Z\s]+(?:\n|$)',  # ALL CAPS TITLE
            r'^(?:Article|ARTICLE)\s+\d+[\.:]\s*(.*?)(?:\n|$)',  # Article X: Title
        ]
        
        for pattern in section_patterns:
            match = re.search(pattern, text, re.MULTILINE)
            if match:
                return match.group(1).strip() if len(match.groups()) > 0 else match.group(0).strip()
        
        # Fallback: use first line if it's short enough to be a title
        lines = text.split('\n')
        if lines and len(lines[0]) < 100:
            return lines[0].strip()
        
        return "Untitled Section"
    
    def generate_chunk_summary(self, text: str, max_retries: int = 6) -> str:
        """
        Generate a summary for a chunk using LLM with production-grade retry logic.
        
        Args:
            text: The chunk text
            max_retries: Maximum number of retry attempts
            
        Returns:
            Generated summary in the format "Label - brief summary sentence"
        """
        import random
        
        for attempt in range(max_retries):
            try:
                prompt = f"""
                Create a concise label that captures the essence of this text chunk, followed by a dash, 
                and then a summary sentence of no more than 12 words.
                
                Format example: "Employment Benefits - Outlines available health insurance and retirement options."
                
                Text chunk:
                {text[:2000]}  # Limit input to avoid token limits
                
                Label - Summary:
                """
                
                response = self.llm.invoke(prompt)
                summary = response.content.strip()
                
                # Ensure it has the label - summary format
                if " - " not in summary:
                    parts = summary.split(" ", 3)
                    if len(parts) >= 3:
                        summary = f"{parts[0]} {parts[1]} - {' '.join(parts[2:])}"
                    else:
                        summary = f"Section Summary - {summary}"
                print(f"LLM attempt {attempt + 1} succeeded") # this is a debug statement
                return summary
                
            except Exception as e:
                print(f"LLM attempt {attempt + 1} failed: {e}")
                if "rate_limit" in str(e).lower() or "429" in str(e):
                    # Exponential backoff with jitter for rate limits
                    base_delay = 2 ** attempt
                    jitter = random.uniform(0.1, 0.3) * base_delay
                    wait_time = base_delay + jitter
                    print(f"Rate limited on attempt {attempt + 1}, waiting {wait_time:.1f}s")
                    time.sleep(wait_time)
                elif "timeout" in str(e).lower():
                    # Shorter delay for timeouts
                    wait_time = 1 + (attempt * 0.5)
                    print(f"Timeout on attempt {attempt + 1}, waiting {wait_time:.1f}s")
                    time.sleep(wait_time)
                else:
                    print(f"LLM error attempt {attempt + 1}: {e}")
                    time.sleep(1)  # Brief pause for other errors
                
                if attempt == max_retries - 1:
                    print(f"All {max_retries} attempts failed, raising exception")
                    raise Exception(f"All {max_retries} LLM attempts failed: {e}")

        print("ERROR: Reached end of generate_chunk_summary without return or exception")

    def extract_keywords(self, text: str) -> List[str]:
        """
        Extract semantic keywords from the chunk text.
        
        Args:
            text: The chunk text
            
        Returns:
            List of extracted keywords
        """
        # Simple rule-based keyword extraction
        # In a production system, you might use NLP techniques or LLMs for this
        words = re.findall(r'\b[a-zA-Z]{4,}\b', text.lower())
        word_counts = {}
        
        for word in words:
            if word not in ['that', 'this', 'with', 'from', 'have', 'were', 'they', 'their']:
                word_counts[word] = word_counts.get(word, 0) + 1
        
        # Get top keywords by frequency
        sorted_words = sorted(word_counts.items(), key=lambda x: x[1], reverse=True)
        return [word for word, _ in sorted_words[:10]]  # Return top 10 keywords
    
    def process_file(self, file_path: Path, document_summary: str = "") -> List[Dict[str, Any]]:
        """
        Process a file into chunks with rich metadata.
        
        Args:
            file_path: Path to the document file
            document_summary: Optional summary for the whole document
            
        Returns:
            List of chunk dictionaries with metadata
        """
        # Calculate file hash (for change detection)
        file_hash = self.calculate_file_hash(file_path)
        file_size = file_path.stat().st_size

        # Extract text from file
        text = self.extract_text_from_file(file_path)
        if not text:
            return []
        
        # Split text into chunks
        chunks = self.text_splitter.split_text(text)
        print(f"Split '{file_path.name}' into {len(chunks)} chunks")
        
        # Process each chunk
        processed_chunks = []
        
        # If no document summary was provided, generate one
        if not document_summary:
            try:
                doc_summary_prompt = f"Provide a concise 8-word summary of this document title: {file_path.stem}"
                document_summary = self.llm.invoke(doc_summary_prompt).content.strip()
            except:
                document_summary = f"{file_path.stem} policy document"
        
        # First pass to create basic chunks
        basic_chunks = []
        for i, chunk_text in enumerate(chunks):
            chunk_id = f"{file_path.stem}_chunk_{i}"
            section_title = self.extract_section_title(chunk_text)
            
            chunk = {
                "id": chunk_id,
                "text": chunk_text,
                "chunk_id": chunk_id,
                "filename": file_path.name,
                "file_type": file_path.suffix.lstrip('.'),
                "file_path": str(file_path),
                "file_hash": file_hash,  
                "file_size": file_size, 
                "indexed_at": datetime.datetime.now().isoformat(),
                "section_title": section_title,
                "document_summary": document_summary,
                "semantic_keywords": self.extract_keywords(chunk_text)
            }
            basic_chunks.append(chunk)
        
        # Second pass to add related chunks and cross-chunk context
        for i, chunk in enumerate(basic_chunks):
            # Find related chunks (previous and next, if they exist)
            related_chunks = []
            if i > 0:
                related_chunks.append(basic_chunks[i-1]["chunk_id"])
            if i < len(basic_chunks) - 1:
                related_chunks.append(basic_chunks[i+1]["chunk_id"])
            
            chunk["related_chunks"] = related_chunks
            
            # Determine section context
            if i > 0 and basic_chunks[i-1]["section_title"] == chunk["section_title"]:
                chunk["section_context"] = f"Section: {chunk['section_title']} (continued)"
            else:
                chunk["section_context"] = f"Section: {chunk['section_title']}"
            
            # Generate chunk summary with production retry logic
            try:
                chunk["chunk_summary"] = self.generate_chunk_summary(chunk["text"])
            except Exception as e:
                print(f"FAILED to generate summary for chunk {i} after all retries: {e}")
                # Since you don't want chunks without summaries, skip this chunk
                continue
            
            chunk["cross_references"] = []
            processed_chunks.append(chunk)
            
            # Rate limiting between chunks (production practice)
            if i < len(basic_chunks) - 1:
                time.sleep(0.2)  # 200ms between chunks
        
        return processed_chunks
    
    def process_directory(self, directory_path: str) -> List[Dict[str, Any]]:
        """
        Process all documents in a directory into chunks.
        
        Args:
            directory_path: Path to directory containing documents
            
        Returns:
            List of chunk dictionaries with metadata
        """
        # Updated supported extensions
        supported_extensions = ['.pdf', '.docx', '.pptx', '.odp', '.txt', '.md']
        
        # Get all files
        directory = Path(directory_path)
        all_files = [p for p in directory.glob('**/*') if p.is_file() and p.suffix.lower() in supported_extensions]
        
        print(f"Found {len(all_files)} supported documents to process in {directory_path}")
        
        all_chunks = []
        
        # Process each file
        for file_path in tqdm(all_files, desc="Processing files"):
            file_chunks = self.process_file(file_path)
            all_chunks.extend(file_chunks)
            print(f"Created {len(file_chunks)} chunks for {file_path.name}")
        
        print(f"Total document chunks created: {len(all_chunks)}")
        return all_chunks
    
    def test_raw_chunk_count(self, directory_path: str) -> int:
        """
        Test method to count raw chunks without any LLM processing.
        
        Args:
            directory_path: Path to directory containing documents
            
        Returns:
            Total number of raw chunks created
        """
        supported_extensions = ['.pdf', '.docx', '.pptx', '.odp', '.txt', '.md']
        
        directory = Path(directory_path)
        all_files = [p for p in directory.glob('**/*') if p.is_file() and p.suffix.lower() in supported_extensions]
        
        total_raw_chunks = 0
        
        print(f"=== RAW CHUNK COUNT TEST ===")
        print(f"Found {len(all_files)} files to process")
        
        for file_path in all_files:
            # Extract text
            text = self.extract_text_from_file(file_path)
            if not text:
                print(f"No text extracted from {file_path.name}")
                continue
                
            # Split into chunks (no processing)
            raw_chunks = self.text_splitter.split_text(text)
            chunk_count = len(raw_chunks)
            total_raw_chunks += chunk_count
            
            print(f"{file_path.name}: {chunk_count} chunks")
        
        print(f"=== TOTAL RAW CHUNKS: {total_raw_chunks} ===")
        return total_raw_chunks
    
def main(cli_args: Optional[List[str]] = None) -> None:
    """CLI entry point for processing a directory of documents."""

    parser = argparse.ArgumentParser(
        description="Process documents in a directory into enriched chunks."
    )
    parser.add_argument(
        "docs_path",
        nargs="?",
        default="Temp_docs_list",
        help="Directory containing documents to process (default: Temp_docs_list)",
    )

    args = parser.parse_args(cli_args)

    openai_client = client_factory.create_openai_client()

    processor = DocumentProcessor(openai_client=openai_client)
    processor.process_directory(args.docs_path)


if __name__ == "__main__":
    main()
